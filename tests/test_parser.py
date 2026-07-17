from pathlib import Path
from types import SimpleNamespace
import sys

import pytest

from engine.models import ParseResult
from engine.parser import parse_file, parse_text


class FakePDFDocument:
    def __init__(self, pages):
        self.pages = pages
        self.page_count = len(pages)
        self.closed = False

    def __iter__(self):
        return iter(self.pages)

    def close(self):
        self.closed = True


class FakePDFPage:
    def __init__(self, plain_text, blocks):
        self.plain_text = plain_text
        self.blocks = blocks

    def get_text(self, mode=None):
        if mode == "blocks":
            return self.blocks
        return self.plain_text


def _install_fake_fitz(monkeypatch, pages):
    document = FakePDFDocument(pages)
    monkeypatch.setitem(sys.modules, "fitz", SimpleNamespace(open=lambda path: document))
    return document


def _org_chart_blocks():
    blocks = [
        (450, 80, 550, 95, "ORG CHART", 0, 0),
        (450, 100, 550, 115, "Nico Reimel", 1, 0),
        (452, 116, 548, 130, "Off Cycle", 2, 0),
    ]
    for index, x_center in enumerate([220, 360, 500, 640, 780], start=3):
        blocks.extend(
            [
                (x_center - 45, 250, x_center + 45, 265, f"Person {index}", index, 0),
                (x_center - 45, 266, x_center + 45, 280, f"Role {index}", index + 20, 0),
            ]
        )
    return blocks


def test_parse_text_returns_manual_parse_result():
    parsed = parse_text("今天面试了一个自动驾驶 CTO 岗位。", source_name="manual_note")

    assert parsed == ParseResult(
        text="今天面试了一个自动驾驶 CTO 岗位。",
        source_name="manual_note",
        source_type="text",
        metadata={"input": "manual"},
    )


async def test_parse_txt_and_markdown_files(tmp_path):
    txt = tmp_path / "note.txt"
    md = tmp_path / "plan.md"
    txt.write_text("纯文本内容", encoding="utf-8")
    md.write_text("# 标题\n\n## 小节\nMarkdown 内容", encoding="utf-8")

    parsed_txt = await parse_file(str(txt))
    parsed_md = await parse_file(str(md))

    assert parsed_txt.text == "纯文本内容"
    assert parsed_txt.source_type == "txt"
    assert parsed_txt.metadata["coverage"]["counts"]["characters"] == 5
    assert parsed_txt.metadata["coverage"]["status"] == "complete"
    assert parsed_txt.quality is not None
    assert parsed_txt.quality.status == "high"
    assert "Markdown 内容" in parsed_md.text
    assert parsed_md.source_type == "md"
    assert parsed_md.quality is not None
    assert parsed_md.quality.status == "high"


async def test_parse_docx_extracts_all_paragraphs(tmp_path):
    docx = pytest.importorskip("docx")
    path = tmp_path / "report.docx"
    document = docx.Document()
    document.add_paragraph("第一段")
    document.add_paragraph("第二段")
    document.add_paragraph("第三段")
    document.save(path)

    parsed = await parse_file(str(path))

    assert parsed.source_type == "docx"
    assert "第一段\n第二段\n第三段" in parsed.text
    assert parsed.metadata["paragraph_count"] == 3
    assert parsed.quality is not None
    assert parsed.quality.status == "high"


async def test_parse_docx_extracts_table_cells_and_reports_coverage(tmp_path):
    docx = pytest.importorskip("docx")
    path = tmp_path / "table-report.docx"
    document = docx.Document()
    document.add_paragraph("人员安排")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "职责"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "项目负责人"
    document.save(path)

    parsed = await parse_file(str(path))

    assert "人员安排" in parsed.text
    assert "| 姓名 | 职责 |" in parsed.text
    assert "| 张三 | 项目负责人 |" in parsed.text
    assert parsed.metadata["coverage"] == {
        "format": "docx",
        "status": "complete",
        "warnings": [],
        "counts": {"paragraphs": 1, "tables": 1, "table_rows": 2},
    }


async def test_parse_pptx_extracts_slide_text(tmp_path):
    pptx = pytest.importorskip("pptx")
    path = tmp_path / "deck.pptx"
    deck = pptx.Presentation()
    for text in ["第一页内容", "第二页内容"]:
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        textbox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
        textbox.text = text
    deck.save(path)

    parsed = await parse_file(str(path))

    assert parsed.source_type == "pptx"
    assert "第一页内容" in parsed.text
    assert "第二页内容" in parsed.text
    assert parsed.metadata["slide_count"] == 2
    assert parsed.quality is not None
    assert parsed.quality.status == "high"


async def test_parse_pptx_extracts_tables_and_notes_with_coverage(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    path = tmp_path / "table-deck.pptx"
    deck = pptx.Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "岗位"
    table.cell(0, 1).text = "人数"
    table.cell(1, 0).text = "工程师"
    table.cell(1, 1).text = "10"
    slide.notes_slide.notes_text_frame.text = "备注：该页人数为规划值"
    deck.save(path)

    parsed = await parse_file(str(path))

    assert "| 岗位 | 人数 |" in parsed.text
    assert "| 工程师 | 10 |" in parsed.text
    assert "备注：该页人数为规划值" in parsed.text
    assert parsed.metadata["coverage"]["counts"] == {"slides": 1, "tables": 1, "notes": 1}
    assert parsed.metadata["coverage"]["status"] == "complete"


async def test_parse_pdf_extracts_all_pages(tmp_path):
    fitz = pytest.importorskip("fitz")
    path = tmp_path / "brief.pdf"
    doc = fitz.open()
    for text in ["第一页 PDF 内容", "第二页 PDF 内容"]:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    doc.save(path)
    doc.close()

    parsed = await parse_file(str(path))

    assert parsed.source_type == "pdf"
    assert "PDF" in parsed.text
    assert "## Page" not in parsed.text
    assert parsed.metadata["page_count"] == 2
    assert parsed.metadata["non_empty_pages"] == 2
    assert parsed.metadata["coverage"]["counts"] == {"pages": 2, "non_empty_pages": 2, "org_chart_pages": 0}
    assert parsed.quality is not None


async def test_parse_pdf_detects_org_chart_page_and_emits_pre_chunk(monkeypatch, tmp_path):
    path = tmp_path / "jlr_org.pdf"
    path.write_bytes(b"%PDF fake")
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "ORG CHART\nNico Reimel\nOff Cycle\nJames Vallance\nConcepts",
                _org_chart_blocks(),
            )
        ],
    )

    parsed = await parse_file(str(path), extract_org_charts=True)

    assert len(parsed.pre_chunks) == 1
    pre_chunk = parsed.pre_chunks[0]
    assert pre_chunk.source_type == "org_chart"
    assert pre_chunk.is_pre_chunked is True
    assert "[ORG_CHART]" in pre_chunk.text
    assert pre_chunk.metadata["page"] == 1
    assert pre_chunk.metadata["org_chart_mode"] == "pdf_layout_fallback"


async def test_parse_pdf_ignores_org_chart_fallback_by_default(monkeypatch, tmp_path):
    path = tmp_path / "jlr_org.pdf"
    path.write_bytes(b"%PDF fake")
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "ORG CHART\nNico Reimel\nOff Cycle\nJames Vallance\nConcepts",
                _org_chart_blocks(),
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "Nico Reimel" in parsed.text


async def test_parse_pdf_detects_org_chart_when_explicitly_enabled(monkeypatch, tmp_path):
    path = tmp_path / "jlr_org.pdf"
    path.write_bytes(b"%PDF fake")
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "ORG CHART\nNico Reimel\nOff Cycle\nJames Vallance\nConcepts",
                _org_chart_blocks(),
            )
        ],
    )

    parsed = await parse_file(str(path), extract_org_charts=True)

    assert len(parsed.pre_chunks) == 1
    assert parsed.pre_chunks[0].source_type == "org_chart"


async def test_parse_pdf_splits_large_org_chart_projection_for_embedding_safety(monkeypatch, tmp_path):
    path = tmp_path / "large_org.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [(450, 80, 550, 95, "ORG CHART", 0, 0)]
    for index in range(90):
        x_center = 150 + (index % 6) * 130
        y = 140 + index * 18
        blocks.extend(
            [
                (x_center - 45, y, x_center + 45, y + 12, f"Person {index:02d}", index * 2 + 1, 0),
                (x_center - 45, y + 13, x_center + 45, y + 25, f"Role {index:02d}", index * 2 + 2, 0),
            ]
        )
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "ORG CHART\n" + "\n".join(f"Person {index:02d}\nRole {index:02d}" for index in range(90)),
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path), extract_org_charts=True)

    assert len(parsed.pre_chunks) > 1
    assert all(record.source_type == "org_chart" for record in parsed.pre_chunks)
    assert all(record.is_pre_chunked is True for record in parsed.pre_chunks)
    assert all(len(record.text) <= 8000 for record in parsed.pre_chunks)
    assert any("[ORG_CHART_SUBTREE]" in record.text for record in parsed.pre_chunks)


async def test_org_chart_page_is_removed_from_normal_pdf_text(monkeypatch, tmp_path):
    path = tmp_path / "mixed.pdf"
    path.write_bytes(b"%PDF fake")
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "ORG CHART\nNico Reimel\nOff Cycle\nJames Vallance\nConcepts",
                _org_chart_blocks(),
            ),
            FakePDFPage(
                "This normal paragraph discusses programme milestones and delivery risks in full sentences.",
                [
                    (
                        72,
                        72,
                        500,
                        96,
                        "This normal paragraph discusses programme milestones and delivery risks in full sentences.",
                        0,
                        0,
                    )
                ],
            ),
        ],
    )

    parsed = await parse_file(str(path), extract_org_charts=True)

    assert "This normal paragraph discusses programme milestones" in parsed.text
    assert "Nico Reimel" not in parsed.text
    assert "James Vallance" not in parsed.text


async def test_non_org_chart_pdf_keeps_existing_parse_behavior(monkeypatch, tmp_path):
    path = tmp_path / "normal.pdf"
    path.write_bytes(b"%PDF fake")
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "This is a normal PDF page with full paragraph text about vehicle programme delivery.",
                [
                    (
                        72,
                        72,
                        500,
                        96,
                        "This is a normal PDF page with full paragraph text about vehicle programme delivery.",
                        0,
                        0,
                    )
                ],
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.source_type == "pdf"
    assert parsed.pre_chunks == []
    assert "normal PDF page" in parsed.text


async def test_parameter_table_pdf_page_is_not_detected_as_org_chart(monkeypatch, tmp_path):
    path = tmp_path / "vision_operator.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [
        (72, 72, 360, 88, "视觉内容理解（Doubao-1.5-vision-pro-32k）", 0, 0),
        (72, 120, 150, 134, "参数名称", 1, 0),
        (180, 120, 230, 134, "类型", 2, 0),
        (260, 120, 330, 134, "默认值", 3, 0),
        (360, 120, 460, 134, "描述", 4, 0),
        (72, 150, 150, 164, "model", 5, 0),
        (180, 150, 230, 164, "str", 6, 0),
        (260, 150, 330, 164, "必填", 7, 0),
        (360, 150, 520, 164, "模型名称", 8, 0),
        (72, 180, 150, 194, "Python", 9, 0),
        (180, 180, 330, 194, "df.with_column(", 10, 0),
        (360, 180, 520, 194, "ArkLLMVisionUnderstanding", 11, 0),
        (72, 210, 150, 224, "images", 12, 0),
        (180, 210, 330, 224, "传入待处理图片", 13, 0),
        (360, 210, 520, 224, "支持单条或多条", 14, 0),
    ]
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "参数名称 类型 默认值 描述\nPython\ndf.with_column(\nArkLLMVisionUnderstanding\nimages 传入待处理图片",
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "参数名称 类型 默认值 描述" in parsed.text


async def test_table_of_contents_pdf_page_is_not_detected_as_org_chart(monkeypatch, tmp_path):
    path = tmp_path / "product_plan.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [
        (72, 72, 360, 88, "0. 文档控制与修订摘要", 0, 0),
        (72, 100, 360, 116, "目录", 1, 0),
        (72, 128, 300, 144, "1. 执行摘要与总判断", 2, 0),
        (330, 128, 380, 144, "2", 3, 0),
        (72, 156, 300, 172, "2. 产品定义与应用边界", 4, 0),
        (330, 156, 380, 172, "5", 5, 0),
        (72, 184, 300, 200, "3. 系统架构与技术方案", 6, 0),
        (330, 184, 380, 200, "8", 7, 0),
        (72, 212, 300, 228, "4. 关键指标、验收方法与数据规范", 8, 0),
        (330, 212, 380, 228, "12", 9, 0),
        (72, 240, 300, 256, "5. 研发可行性、风险与修正方案", 10, 0),
        (330, 240, 380, 256, "16", 11, 0),
        (72, 268, 300, 284, "附录 A. V2 问题到 V3 修正对照", 12, 0),
        (330, 268, 380, 284, "22", 13, 0),
    ]
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "目录\n1. 执行摘要与总判断\n2. 产品定义与应用边界\n3. 系统架构与技术方案\n附录 A. V2 问题到 V3 修正对照",
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "目录" in parsed.text


async def test_itinerary_table_pdf_page_is_not_detected_as_org_chart(monkeypatch, tmp_path):
    path = tmp_path / "itinerary.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [
        (72, 72, 260, 88, "二、 每日行程概要", 0, 0),
        (72, 120, 140, 134, "日期", 1, 0),
        (150, 120, 260, 134, "核心路线节点", 2, 0),
        (270, 120, 330, 134, "里程 (km)", 3, 0),
        (340, 120, 430, 134, "驾车时长", 4, 0),
        (440, 120, 560, 134, "今日行程亮点", 5, 0),
        (72, 150, 140, 164, "6月7日", 6, 0),
        (150, 150, 260, 164, "北京 ➔ 伊宁机场", 7, 0),
        (270, 150, 330, 164, "5 km", 8, 0),
        (340, 150, 430, 164, "15 min", 9, 0),
        (440, 150, 560, 164, "六星街烤肉", 10, 0),
        (72, 180, 140, 194, "6月8日", 11, 0),
        (150, 180, 260, 194, "赛里木湖环湖", 12, 0),
        (270, 180, 330, 194, "160 km", 13, 0),
        (340, 180, 430, 194, "2.5 h", 14, 0),
        (440, 180, 560, 194, "雪山与草原", 15, 0),
    ]
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "二、 每日行程概要\n日期 核心路线节点 里程 (km) 驾车时长 今日行程亮点\n"
                "6月7日 北京 ➔ 伊宁机场 5 km 15 min 六星街烤肉\n"
                "6月8日 赛里木湖环湖 160 km 2.5 h 雪山与草原",
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "每日行程概要" in parsed.text


async def test_travel_preparation_page_is_not_detected_as_org_chart(monkeypatch, tmp_path):
    path = tmp_path / "travel_prep.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [
        (72, 72, 360, 88, "新疆伊犁河谷与独库公路自驾终极路书", 0, 0),
        (72, 110, 220, 124, "出行时间：2026年6月7日 - 6月12日", 1, 0),
        (250, 110, 360, 124, "出行人数：2人核心团队", 2, 0),
        (72, 140, 520, 154, "核心路线：伊宁-赛湖-独山子大峡谷-独库北段-唐布拉-伊宁", 3, 0),
        (72, 180, 260, 194, "一、 行前准备（关键要点）", 4, 0),
        (90, 210, 520, 224, "• 离线地图：务必提前下载伊犁州、博州、塔城地区离线包。", 5, 0),
        (90, 240, 520, 254, "• 衣物储备：准备轻薄羽绒服、冲锋衣和中高帮徒步鞋。", 6, 0),
        (90, 270, 520, 284, "• 随车物品：手机支架、快充头、移动电源、湿巾纸巾。", 7, 0),
        (90, 300, 520, 314, "• 预约状态明晰：独库公路国道无需预约。", 8, 0),
    ]
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "新疆伊犁河谷与独库公路自驾终极路书\n"
                "出行时间：2026年6月7日 - 6月12日\n"
                "出行人数：2人核心团队\n"
                "核心路线：伊宁-赛湖-独山子大峡谷-独库北段-唐布拉-伊宁\n"
                "一、 行前准备（关键要点）\n"
                "• 离线地图：务必提前下载伊犁州、博州、塔城地区离线包。\n"
                "• 衣物储备：准备轻薄羽绒服、冲锋衣和中高帮徒步鞋。\n"
                "• 随车物品：手机支架、快充头、移动电源、湿巾纸巾。\n"
                "• 预约状态明晰：独库公路国道无需预约。",
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "行前准备" in parsed.text


async def test_milestone_table_pdf_page_is_not_detected_as_org_chart(monkeypatch, tmp_path):
    path = tmp_path / "milestones.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [
        (72, 72, 300, 88, "6. 开发计划与阶段里程碑", 0, 0),
        (72, 120, 120, 134, "阶段", 1, 0),
        (130, 120, 190, 134, "周期", 2, 0),
        (200, 120, 300, 134, "目标", 3, 0),
        (310, 120, 420, 134, "关键交付物", 4, 0),
        (430, 120, 560, 134, "Go/No-Go 条件", 5, 0),
        (72, 150, 120, 164, "需求冻结/RFQ", 6, 0),
        (130, 150, 190, 164, "第 0-2 周", 7, 0),
        (200, 150, 300, 164, "冻结场景和数据 schema", 8, 0),
        (310, 150, 420, 164, "PRD、RFQ 包", 9, 0),
        (430, 150, 560, 164, "接口满足触发要求", 10, 0),
        (72, 180, 120, 194, "EVT 样机", 11, 0),
        (130, 180, 190, 194, "第 2-8 周", 12, 0),
        (200, 180, 300, 194, "完成 3-5 套工程样机", 13, 0),
        (310, 180, 420, 194, "BOM Rev.A", 14, 0),
        (430, 180, 560, 194, "P95 同步通过", 15, 0),
    ]
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "6. 开发计划与阶段里程碑\n阶段 周期 目标 关键交付物 Go/No-Go 条件\n"
                "需求冻结/RFQ 第 0-2 周 冻结场景和数据 schema PRD、RFQ 包 接口满足触发要求\n"
                "EVT 样机 第 2-8 周 完成 3-5 套工程样机 BOM Rev.A P95 同步通过",
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "开发计划与阶段里程碑" in parsed.text


async def test_automotive_business_diagram_pdf_page_is_not_detected_as_org_chart(monkeypatch, tmp_path):
    path = tmp_path / "honda_business.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [
        (72, 72, 360, 88, "Advancement in the application of intelligent technologies", 0, 0),
        (72, 110, 180, 124, "Next-generation ADAS", 1, 0),
        (210, 110, 360, 124, "Expressway ADAS", 2, 0),
        (390, 110, 540, 124, "Surface road ADAS", 3, 0),
        (72, 145, 170, 159, "Autonomous parking assist", 4, 0),
        (210, 145, 360, 159, "NO TURN ON RED STOP", 5, 0),
        (390, 145, 540, 159, "Shopping Office Travel", 6, 0),
        (72, 190, 220, 204, "Advancing e:HEV and platforms", 7, 0),
        (250, 190, 360, 204, "Fuel economy", 8, 0),
        (390, 190, 470, 204, "10%", 9, 0),
        (72, 230, 220, 244, "Realignment of EV strategy", 10, 0),
        (250, 230, 360, 244, "Carbon neutrality", 11, 0),
        (390, 230, 470, 244, "2030 2035 2040", 12, 0),
        (72, 270, 220, 284, "Motorcycle business", 13, 0),
        (250, 270, 360, 284, "Global market share", 14, 0),
        (390, 270, 470, 284, "ROIC 10％", 15, 0),
    ]
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "Advancement in the application of intelligent technologies\n"
                "Next-generation ADAS Expressway ADAS Surface road ADAS Autonomous parking assist\n"
                "Advancing e:HEV and platforms Fuel economy Improvement by more than 10%\n"
                "Realignment of EV strategy Carbon neutrality 2030 2035 2040\n"
                "Motorcycle business Global market share Operating profit ROIC 10％",
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "Next-generation ADAS" in parsed.text


async def test_platform_deployment_architecture_pdf_page_is_not_detected_as_org_chart(monkeypatch, tmp_path):
    path = tmp_path / "celonis_omnia.pdf"
    path.write_bytes(b"%PDF fake")
    blocks = [
        (72, 72, 360, 88, "Celonis “Omnia” - self-managed deployment", 0, 0),
        (72, 110, 280, 124, "partner- or customer-managed Kubernetes clusters", 1, 0),
        (310, 110, 480, 124, "dedicated OCI registry", 2, 0),
        (72, 150, 190, 164, "AWS or Azure Cloud Account", 3, 0),
        (220, 150, 360, 164, "Customer Cloud Account", 4, 0),
        (390, 150, 520, 164, "Customer Datacenter", 5, 0),
        (72, 190, 190, 204, "Storage & Compute", 6, 0),
        (220, 190, 360, 204, "Shared multi-tenant realm", 7, 0),
        (390, 190, 520, 204, "Single-tenant realm", 8, 0),
        (72, 230, 190, 244, "AI-Powered Process Copilots", 9, 0),
        (220, 230, 360, 244, "AI Apps", 10, 0),
        (390, 230, 520, 244, "AI for Data Enrichment", 11, 0),
        (72, 270, 190, 284, "ML Workbench", 12, 0),
        (220, 270, 360, 284, "AI Integrations & Partnerships", 13, 0),
    ]
    _install_fake_fitz(
        monkeypatch,
        [
            FakePDFPage(
                "Celonis “Omnia” - self-managed deployment\n"
                "Enable flexible deployment on partner- or customer-managed Kubernetes clusters\n"
                "Distribute Celonis as releases via a dedicated OCI registry\n"
                "AWS or Azure Cloud Account Customer Cloud Account Customer Datacenter\n"
                "Storage & Compute Shared multi-tenant realm Single-tenant realm\n"
                "AI-Powered Process Copilots AI Apps AI for Data Enrichment ML Workbench",
                blocks,
            )
        ],
    )

    parsed = await parse_file(str(path))

    assert parsed.pre_chunks == []
    assert "self-managed deployment" in parsed.text


async def test_parse_pdf_cleaning_reassesses_quality(tmp_path):
    fitz = pytest.importorskip("fitz")
    path = tmp_path / "paged.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Page 1\n智能座舱市场规模持续增长，2026 年预计达到 1200 亿元。")
    doc.save(path)
    doc.close()

    parsed = await parse_file(str(path))

    assert "Page 1" not in parsed.text
    assert parsed.quality is not None
    assert parsed.quality.status in {"high", "low"}
    assert parsed.metadata["quality_status"] == parsed.quality.status


async def test_parse_xlsx_converts_sheets_to_markdown_tables(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    path = tmp_path / "data.xlsx"
    workbook = openpyxl.Workbook()
    first = workbook.active
    first.title = "候选人"
    first.append(["姓名", "评分"])
    first.append(["TZ", 95])
    second = workbook.create_sheet("项目")
    second.append(["项目", "状态"])
    second.append(["PKA", "进行中"])
    workbook.save(path)

    parsed = await parse_file(str(path))

    assert parsed.source_type == "xlsx"
    assert "## Sheet: 候选人" in parsed.text
    assert "| 姓名 | 评分 |" in parsed.text
    assert "| PKA | 进行中 |" in parsed.text
    assert parsed.metadata["sheet_count"] == 2
    assert parsed.quality is not None
    assert parsed.quality.status == "high"


async def test_parse_xlsx_preserves_formula_expression_and_reports_coverage(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    path = tmp_path / "formula.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "预算"
    sheet.append(["项目", "金额"])
    sheet.append(["合计", "=SUM(10,20)"])
    workbook.save(path)

    parsed = await parse_file(str(path))

    assert "=SUM(10,20)" in parsed.text
    assert parsed.metadata["coverage"]["counts"] == {"sheets": 1, "rows": 2, "formulas": 1}
    assert parsed.metadata["coverage"]["status"] == "complete"


class FakeOCR:
    async def extract(self, image_paths):
        assert image_paths == [self.expected_path]
        return "图片里的中文"


class UsefulImageOCR:
    async def extract(self, image_paths):
        return "\n".join(
            [
                "核心一句话总结 Executive One-liner",
                "通过品牌矩阵、动力多元化、北美聚焦和成本优化，JLR 进入 Reimagine 战略增长兑现阶段。",
                "Growth Reimagined strategy delivery phase and portfolio strategy.",
                "Range Rover Defender Discovery Jaguar market focus North America.",
            ]
        )


class FailingOCR:
    async def extract(self, image_paths):
        raise RuntimeError("OCR failed after 3 retries: timeout")


async def test_parse_image_uses_injected_ocr_client_inside_running_event_loop(tmp_path):
    path = tmp_path / "screenshot.png"
    path.write_bytes(b"\x89PNG\r\n\x1a\n")
    ocr = FakeOCR()
    ocr.expected_path = str(path)

    parsed = await parse_file(str(path), ocr_client=ocr)

    assert parsed.source_type == "image"
    assert parsed.text == "图片里的中文"
    assert parsed.metadata["ocr"] is True


async def test_parse_image_marks_useful_ocr_quality(tmp_path):
    path = tmp_path / "screenshot.jpeg"
    path.write_bytes(b"jpeg bytes")

    parsed = await parse_file(str(path), mime_type="image/jpeg", ocr_client=UsefulImageOCR())

    assert parsed.source_type == "image"
    assert parsed.quality is not None
    assert parsed.quality.status == "high"
    assert parsed.quality.action == "image_ocr"
    assert parsed.quality.effective_chars_per_page >= 80


async def test_parse_image_marks_short_ocr_as_low_quality(tmp_path):
    path = tmp_path / "tiny.jpeg"
    path.write_bytes(b"jpeg bytes")
    ocr = FakeOCR()
    ocr.expected_path = str(path)

    parsed = await parse_file(str(path), mime_type="image/jpeg", ocr_client=ocr)

    assert parsed.source_type == "image"
    assert parsed.quality is not None
    assert parsed.quality.status == "low"
    assert parsed.quality.action == "image_ocr_low"
    assert any("OCR 文本少于 80 字" in reason for reason in parsed.quality.reasons)


async def test_parse_image_raises_clear_error_when_ocr_fails(tmp_path):
    path = tmp_path / "broken.png"
    path.write_bytes(b"not a real image")

    with pytest.raises(RuntimeError, match="OCR failed after 3 retries"):
        await parse_file(str(path), ocr_client=FailingOCR())


class EmptyOCR:
    async def extract(self, image_paths):
        return "   \n  "


async def test_parse_image_rejects_empty_ocr_text(tmp_path):
    path = tmp_path / "photo.jpeg"
    path.write_bytes(b"jpeg bytes")

    with pytest.raises(ValueError, match="OCR produced no usable text for image"):
        await parse_file(str(path), mime_type="image/jpeg", ocr_client=EmptyOCR())


async def test_parse_corrupt_docx_raises_clear_error(tmp_path):
    path = tmp_path / "broken.docx"
    path.write_text("not a real docx", encoding="utf-8")

    with pytest.raises(ValueError, match="Failed to parse"):
        await parse_file(str(path))
