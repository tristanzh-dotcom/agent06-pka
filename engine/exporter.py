from pathlib import Path
from typing import Dict, List

from pptx.util import Pt


def export_to_word(question: str, answer: str, sources: List[Dict], output_path: str) -> str:
    import docx

    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)

    document = docx.Document()
    document.add_heading("PKA 问答导出", level=1)
    document.add_heading("问题", level=2)
    document.add_paragraph(question)
    document.add_heading("回答", level=2)
    document.add_paragraph(answer)
    document.add_heading("参考来源", level=2)

    table = document.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    headers = ["来源文件", "Chunk", "相关度", "Chunk ID"]
    for cell, value in zip(table.rows[0].cells, headers):
        cell.text = value

    for source in sources:
        row = table.add_row().cells
        row[0].text = str(source.get("source_name", ""))
        row[1].text = str(source.get("chunk_index", ""))
        row[2].text = _format_relevance(source.get("relevance"))
        row[3].text = str(source.get("chunk_id", ""))

    document.save(path)
    return str(path)


def export_to_ppt(question: str, answer: str, sources: List[Dict], output_path: str) -> str:
    import pptx

    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)

    presentation = pptx.Presentation()
    _add_title_slide(presentation, "PKA 问答导出", question)
    _add_body_slide(presentation, "回答", answer)
    _add_sources_slide(presentation, sources)
    presentation.save(path)
    return str(path)


def _add_title_slide(presentation, title: str, subtitle: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_body_slide(presentation, title: str, body: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for index, paragraph_text in enumerate(_split_slide_paragraphs(body)):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = paragraph_text
        paragraph.font.size = Pt(18)


def _add_sources_slide(presentation, sources: List[Dict]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "参考来源"
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    source_lines = _source_lines(sources)
    for index, line in enumerate(source_lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(14)


def _split_slide_paragraphs(text: str) -> List[str]:
    paragraphs = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    if not paragraphs and text:
        paragraphs = [str(text).strip()]
    if not paragraphs:
        return ["无回答内容"]
    return [_truncate(paragraph, 320) for paragraph in paragraphs[:8]]


def _source_lines(sources: List[Dict]) -> List[str]:
    if not sources:
        return ["无"]
    lines = []
    for source in sources[:12]:
        source_name = str(source.get("source_name", "")).strip() or "未知来源"
        chunk = str(source.get("chunk_index", "")).strip()
        chunk_id = str(source.get("chunk_id", "")).strip()
        relevance = _format_relevance(source.get("relevance"))
        parts = [source_name]
        if chunk:
            parts.append(f"chunk {chunk}")
        if relevance:
            parts.append(f"相关度 {relevance}")
        if chunk_id:
            parts.append(chunk_id)
        lines.append(" | ".join(parts))
    return lines


def _truncate(text: str, limit: int) -> str:
    return text if len(text) <= limit else text[:limit].rstrip() + "..."


def _format_relevance(value) -> str:
    if value is None:
        return ""
    try:
        return f"{float(value):.4f}"
    except (TypeError, ValueError):
        return str(value)
