import mimetypes
from pathlib import Path
import re
from typing import Any, Optional

from engine.models import ParseResult, PreChunkedParseRecord
from engine.quality import assess_extracted_text_quality
from engine.org_chart import (
    PdfTextBlock,
    generate_projection_text,
    infer_layout_hierarchy,
    merge_pdf_blocks,
    select_org_chart_title,
)


TEXT_TYPES = {".txt": "txt", ".md": "md"}
IMAGE_TYPES = {".png", ".jpg", ".jpeg", ".webp"}
ORG_CHART_MAX_PRE_CHUNK_CHARS = 3500
OCR_ORG_CHART_MAX_LINES = 80


def parse_text(text: str, source_name: str = "manual_input") -> ParseResult:
    return ParseResult(
        text=text,
        source_name=source_name,
        source_type="text",
        metadata={"input": "manual"},
    )


async def parse_file(
    file_path: str,
    mime_type: Optional[str] = None,
    ocr_client: Any = None,
    extract_org_charts: bool = False,
) -> ParseResult:
    path = Path(file_path)
    suffix = path.suffix.lower()
    detected_mime = mime_type or mimetypes.guess_type(str(path))[0] or ""

    try:
        if suffix in TEXT_TYPES:
            text = path.read_text(encoding="utf-8")
            return ParseResult(
                text=text,
                source_name=path.name,
                source_type=TEXT_TYPES[suffix],
                metadata={
                    "mime_type": detected_mime,
                    "coverage": _coverage(
                        TEXT_TYPES[suffix],
                        characters=len(text),
                        lines=len(text.splitlines()),
                    ),
                },
                quality=assess_extracted_text_quality(text),
            )
        if suffix == ".docx":
            return _parse_docx(path)
        if suffix == ".pptx":
            return _parse_pptx(path)
        if suffix == ".pdf":
            return _parse_pdf(path, extract_org_charts=extract_org_charts)
        if suffix == ".xlsx":
            return _parse_xlsx(path)
        if suffix in IMAGE_TYPES or detected_mime.startswith("image/"):
            from engine.quality import assess_image_ocr_quality

            if ocr_client is None:
                raise ValueError("OCR client is required for image parsing")
            text = await ocr_client.extract([str(path)])
            if not str(text or "").strip():
                raise ValueError("OCR produced no usable text for image")
            return ParseResult(
                text=text,
                source_name=path.name,
                source_type="image",
                metadata={
                    "ocr": True,
                    "mime_type": detected_mime,
                    "coverage": _coverage("image", images=1, extracted_characters=len(str(text))),
                },
                quality=assess_image_ocr_quality(text),
            )
    except ValueError:
        raise
    except RuntimeError:
        raise
    except Exception as exc:
        raise ValueError(f"Failed to parse {path.name}: {exc}") from exc

    raise ValueError(f"Unsupported file type: {path.suffix or detected_mime}")


def _parse_docx(path: Path) -> ParseResult:
    import docx
    from docx.table import Table

    document = docx.Document(path)
    sections = []
    paragraph_count = 0
    table_count = 0
    table_row_count = 0
    for block in document.iter_inner_content():
        if isinstance(block, Table):
            rows = [[cell.text.strip() for cell in row.cells] for row in block.rows]
            if rows:
                sections.append(_markdown_table(rows))
                table_count += 1
                table_row_count += len(rows)
            continue
        text = str(getattr(block, "text", "") or "").strip()
        if text:
            sections.append(text)
            paragraph_count += 1
    extracted_text = "\n".join(sections)
    return ParseResult(
        text=extracted_text,
        source_name=path.name,
        source_type="docx",
        metadata={
            "paragraph_count": paragraph_count,
            "coverage": _coverage(
                "docx",
                paragraphs=paragraph_count,
                tables=table_count,
                table_rows=table_row_count,
            ),
        },
        quality=assess_extracted_text_quality(extracted_text),
    )


def _parse_pptx(path: Path) -> ParseResult:
    import pptx

    presentation = pptx.Presentation(path)
    texts = []
    table_count = 0
    note_count = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                if rows:
                    slide_texts.append(_markdown_table(rows))
                    table_count += 1
            elif hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        notes_text = str(getattr(slide.notes_slide.notes_text_frame, "text", "") or "").strip()
        if notes_text:
            slide_texts.append("### Notes\n" + notes_text)
            note_count += 1
        if slide_texts:
            texts.append(f"## Slide {slide_number}\n" + "\n".join(slide_texts))
    extracted_text = "\n\n".join(texts)
    return ParseResult(
        text=extracted_text,
        source_name=path.name,
        source_type="pptx",
        metadata={
            "slide_count": len(presentation.slides),
            "coverage": _coverage(
                "pptx",
                slides=len(presentation.slides),
                tables=table_count,
                notes=note_count,
            ),
        },
        quality=assess_extracted_text_quality(extracted_text),
    )


def _parse_pdf(path: Path, extract_org_charts: bool = False) -> ParseResult:
    import fitz
    from engine.quality import assess_pdf_quality, clean_pdf_text

    document = fitz.open(path)
    pages = []
    pre_chunks = []
    try:
        for page_number, page in enumerate(document, start=1):
            page_text = page.get_text().strip()
            raw_blocks = page.get_text("blocks")
            if page_text:
                blocks = _pdf_text_blocks(raw_blocks, page_number)
                if extract_org_charts and _detect_org_chart_page(page_text, blocks):
                    pre_chunks.extend(
                        _org_chart_pre_chunks(
                            source_name=path.name,
                            page_number=page_number,
                            page_text=page_text,
                            blocks=blocks,
                            page_height=_page_height(page, blocks),
                        )
                    )
                    continue
                pages.append(page_text)
        page_count = document.page_count
    finally:
        document.close()
    raw_text = "\n\n".join(pages)
    cleaned_text = clean_pdf_text(raw_text, page_texts=pages, page_count=page_count)
    quality = assess_pdf_quality(raw_text, cleaned_text, page_count, len(pages))
    return ParseResult(
        text=cleaned_text,
        source_name=path.name,
        source_type="pdf",
        metadata={
            "page_count": page_count,
            "non_empty_pages": len(pages),
            "quality_status": quality.status,
            "quality_action": quality.action,
            "org_chart_pages": [record.metadata["page"] for record in pre_chunks],
            "org_chart_chunks": len(pre_chunks),
            "org_chart_mode": "pdf_layout_fallback" if pre_chunks else "",
            "coverage": _coverage(
                "pdf",
                pages=page_count,
                non_empty_pages=len(pages),
                org_chart_pages=len({record.metadata["page"] for record in pre_chunks}),
            ),
        },
        quality=quality,
        pre_chunks=pre_chunks,
    )


def _parse_xlsx(path: Path) -> ParseResult:
    import openpyxl

    workbook = openpyxl.load_workbook(path, data_only=False)
    value_workbook = openpyxl.load_workbook(path, data_only=True)
    sections = []
    total_rows = 0
    formula_count = 0
    for sheet, value_sheet in zip(workbook.worksheets, value_workbook.worksheets):
        rows = []
        for formula_row, value_row in zip(sheet.iter_rows(values_only=False), value_sheet.iter_rows(values_only=False)):
            rendered = []
            for formula_cell, value_cell in zip(formula_row, value_row):
                if formula_cell.data_type == "f":
                    formula_count += 1
                    formula = str(formula_cell.value or "")
                    formula = formula if formula.startswith("=") else "=" + formula
                    cached = _cell_to_text(value_cell.value)
                    rendered.append(f"{formula}（计算值：{cached}）" if cached else formula)
                else:
                    rendered.append(_cell_to_text(formula_cell.value))
            rows.append(rendered)
        rows = [row for row in rows if any(cell for cell in row)]
        if not rows:
            continue
        total_rows += len(rows)
        width = max(len(row) for row in rows)
        normalized = [row + [""] * (width - len(row)) for row in rows]
        sections.append(f"## Sheet: {sheet.title}\n" + _markdown_table(normalized))
    try:
        extracted_text = "\n\n".join(sections)
        return ParseResult(
            text=extracted_text,
            source_name=path.name,
            source_type="xlsx",
            metadata={
                "sheet_count": len(workbook.worksheets),
                "coverage": _coverage(
                    "xlsx",
                    sheets=len(workbook.worksheets),
                    rows=total_rows,
                    formulas=formula_count,
                ),
            },
            quality=assess_extracted_text_quality(extracted_text),
        )
    finally:
        workbook.close()
        value_workbook.close()


def _cell_to_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


def _coverage(format_name: str, **counts: int) -> dict:
    return {
        "format": format_name,
        "status": "complete",
        "warnings": [],
        "counts": counts,
    }


def _markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    escaped = [[str(cell).replace("|", "\\|").replace("\n", " ") for cell in row] for row in normalized]
    table = [
        "| " + " | ".join(escaped[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    table.extend("| " + " | ".join(row) + " |" for row in escaped[1:])
    return "\n".join(table)


def _pdf_text_blocks(raw_blocks, page_number: int) -> list[PdfTextBlock]:
    blocks = []
    for raw in raw_blocks or []:
        if len(raw) < 5:
            continue
        x0, y0, x1, y1, text = raw[:5]
        cleaned = str(text).strip()
        if not cleaned:
            continue
        blocks.append(
            PdfTextBlock(
                text=cleaned,
                x0=float(x0),
                y0=float(y0),
                x1=float(x1),
                y1=float(y1),
                font_size=max(1.0, min(18.0, float(y1) - float(y0))),
                page=page_number,
            )
        )
    return blocks


def _page_height(page, blocks: list[PdfTextBlock]) -> float:
    rect = getattr(page, "rect", None)
    height = getattr(rect, "height", None)
    if height:
        return float(height)
    return max((block.y1 for block in blocks), default=0.0)


def _detect_org_chart_page(page_text: str, blocks: list[PdfTextBlock]) -> bool:
    normalized = page_text.upper()
    if re.search(r"\bORG(?:ANISATION|ANIZATION)?\s+CHART\b", normalized):
        return True
    if _looks_like_document_table_or_code_page(page_text, blocks):
        return False
    short_blocks = [block for block in blocks if len(block.text) <= 32]
    if len(short_blocks) < 12:
        return False
    y_bands = _count_y_bands(short_blocks)
    x_centers = {round(block.x_center / 80) for block in short_blocks}
    short_ratio = len(short_blocks) / max(len(blocks), 1)
    return short_ratio >= 0.65 and y_bands >= 3 and len(x_centers) >= 3


def _looks_like_document_table_or_code_page(page_text: str, blocks: list[PdfTextBlock]) -> bool:
    text = page_text.strip()
    normalized = text.upper()
    compact_lines = [line.strip() for line in text.splitlines() if line.strip()]
    block_texts = [block.text.strip() for block in blocks if block.text.strip()]
    combined_blocks = "\n".join(block_texts)

    table_markers = ["参数名称", "默认值", "描述", "类型", "输入列名", "输出", "说明"]
    if sum(1 for marker in table_markers if marker in text or marker in combined_blocks) >= 3:
        return True

    itinerary_markers = ["每日行程", "日期", "核心路线", "里程", "驾车时长", "行程亮点", "机场", "酒店"]
    if sum(1 for marker in itinerary_markers if marker in text or marker in combined_blocks) >= 4:
        return True

    travel_prep_markers = ["出行时间", "出行人数", "核心路线", "行前准备", "离线地图", "衣物储备", "随车物品", "预约状态"]
    if sum(1 for marker in travel_prep_markers if marker in text or marker in combined_blocks) >= 4:
        return True

    milestone_markers = ["开发计划", "阶段", "里程碑", "周期", "目标", "关键交付物", "GO/NO-GO", "BOM", "RFQ"]
    if sum(1 for marker in milestone_markers if marker in text or marker in normalized) >= 4:
        return True

    code_markers = ["PYTHON", "DF.WITH_COLUMN", "IMPORT ", "RETURN ", "COL(", "ARKLLMVISIONUNDERSTANDING"]
    if sum(1 for marker in code_markers if marker in normalized) >= 2:
        return True

    toc_markers = ["目录", "文档控制", "修订摘要", "执行摘要", "附录", "BOM", "验收方法"]
    numbered_lines = sum(1 for line in compact_lines if re.match(r"^(?:\d+\.|附录\s*[A-ZＡ-Ｚ]?\b)", line))
    if ("目录" in text or "TABLE OF CONTENTS" in normalized) and numbered_lines >= 3:
        return True
    if sum(1 for marker in toc_markers if marker in text or marker in normalized) >= 3 and numbered_lines >= 2:
        return True

    bullet_or_table_rows = sum(
        1
        for line in compact_lines
        if line.startswith(("•", "◦", "▪", "-", "")) or re.match(r"^[A-Za-z0-9_]+\s+(?:str|int|float|bool|list|dict)\b", line)
    )
    if bullet_or_table_rows >= 5 and any(marker in text for marker in ("参数", "默认", "说明", "输入", "输出")):
        return True

    return False


def _count_y_bands(blocks: list[PdfTextBlock], tolerance: float = 10.0) -> int:
    bands: list[float] = []
    for block in sorted(blocks, key=lambda item: item.y_center):
        for index, center in enumerate(bands):
            if abs(block.y_center - center) <= tolerance:
                bands[index] = (center + block.y_center) / 2
                break
        else:
            bands.append(block.y_center)
    return len(bands)


def _org_chart_pre_chunks(
    *,
    source_name: str,
    page_number: int,
    page_text: str,
    blocks: list[PdfTextBlock],
    page_height: float,
) -> list[PreChunkedParseRecord]:
    title, cleaned_blocks = select_org_chart_title(blocks, page_height)
    candidate_blocks = cleaned_blocks or blocks
    nodes = merge_pdf_blocks(candidate_blocks)
    edges = infer_layout_hierarchy(nodes)
    projection = generate_projection_text(
        source_name=source_name,
        source_page=page_number,
        title=title,
        extraction_mode="pdf_layout_fallback",
        confidence="medium",
        nodes=nodes,
        edges=edges,
        warnings=[
            "native_pptx_unavailable",
            "connector_relationships_inferred",
            "cross_page_links_not_supported_v1",
        ],
    )
    if len(projection) <= ORG_CHART_MAX_PRE_CHUNK_CHARS:
        return [
            _pre_chunk_record(
                text=projection,
                source_name=source_name,
                page_number=page_number,
                part_index=1,
            )
        ]
    return [
        _pre_chunk_record(
            text=text,
            source_name=source_name,
            page_number=page_number,
            part_index=index,
        )
        for index, text in enumerate(
            _split_large_org_chart_projection(
                projection,
                source_name=source_name,
                page_number=page_number,
                title=title,
            ),
            start=1,
        )
    ]


def ocr_org_chart_pre_chunks(
    text: str,
    *,
    source_name: str,
    page_number: int = 1,
) -> list[PreChunkedParseRecord]:
    if not _detect_ocr_org_chart_text(text):
        return []
    lines = _ocr_org_chart_lines(text)
    if not lines:
        return []
    projection = _ocr_org_chart_projection(source_name, page_number, lines)
    return [
        PreChunkedParseRecord(
            text=projection,
            source_name=source_name,
            source_type="org_chart",
            is_pre_chunked=True,
            metadata={
                "page": page_number,
                "chart_id": f"{source_name}#page_{page_number}#ocr_chart_1",
                "confidence": "low",
                "org_chart_mode": "ocr_layout_fallback",
            },
        )
    ]


def _detect_ocr_org_chart_text(text: str) -> bool:
    compact = str(text or "").strip()
    if not compact:
        return False
    normalized = compact.upper()
    has_org_chart_marker = bool(re.search(r"\bORG(?:ANISATION|ANIZATION)?\s+CHART\b", normalized)) or "组织架构" in compact
    if not has_org_chart_marker:
        return False
    lines = _ocr_org_chart_lines(compact)
    short_lines = [line for line in lines if len(line) <= 48]
    return len(short_lines) >= 4


def _ocr_org_chart_lines(text: str) -> list[str]:
    lines = []
    for raw_line in str(text or "").splitlines():
        line = raw_line.strip(" \t-•·|")
        if not line:
            continue
        if re.search(r"\bORG(?:ANISATION|ANIZATION)?\s+CHART\b", line, re.IGNORECASE) or line == "组织架构":
            continue
        lines.append(line)
        if len(lines) >= OCR_ORG_CHART_MAX_LINES:
            break
    return lines


def _ocr_org_chart_projection(source_name: str, page_number: int, lines: list[str]) -> str:
    body = [f"- {line}" for line in lines]
    triggers = [f"- {line} appears in OCR org chart text." for line in lines[:20]]
    return "\n".join(
        [
            "[ORG_CHART_OCR]",
            f"Source: {source_name}",
            f"Page: {page_number}",
            "Extraction mode: ocr_layout_fallback",
            "Confidence: low",
            "",
            "Structure:",
            *body,
            "",
            "Semantic Search Triggers:",
            *triggers,
            "[/ORG_CHART_OCR]",
        ]
    )


def _pre_chunk_record(
    *, text: str, source_name: str, page_number: int, part_index: int
) -> PreChunkedParseRecord:
    return PreChunkedParseRecord(
        text=text,
        source_name=source_name,
        source_type="org_chart",
        is_pre_chunked=True,
        metadata={
            "page": page_number,
            "chart_id": f"{source_name}#page_{page_number}#chart_1_part_{part_index}",
            "confidence": "medium",
            "org_chart_mode": "pdf_layout_fallback",
        },
    )


def _split_large_org_chart_projection(
    projection: str, *, source_name: str, page_number: int, title: str
) -> list[str]:
    body_lines = [
        line
        for line in projection.splitlines()
        if line
        and not line.startswith("[ORG_CHART]")
        and not line.startswith("[/ORG_CHART]")
        and not line.startswith("Source:")
        and not line.startswith("Page:")
        and not line.startswith("Title:")
        and not line.startswith("Extraction mode:")
        and not line.startswith("Confidence:")
    ]
    chunks: list[str] = []
    current: list[str] = []
    for line in body_lines:
        candidate = current + [line]
        if current and len(_org_chart_subtree_text(source_name, page_number, title, candidate)) > ORG_CHART_MAX_PRE_CHUNK_CHARS:
            chunks.append(_org_chart_subtree_text(source_name, page_number, title, current))
            current = [line]
        else:
            current = candidate
    if current:
        chunks.append(_org_chart_subtree_text(source_name, page_number, title, current))
    return chunks


def _org_chart_subtree_text(
    source_name: str, page_number: int, title: str, body_lines: list[str]
) -> str:
    return "\n".join(
        [
            "[ORG_CHART_SUBTREE]",
            f"Source: {source_name}",
            f"Page: {page_number}",
            f"Context Root: {title}",
            "Confidence: medium",
            "",
            *body_lines,
            "[/ORG_CHART_SUBTREE]",
        ]
    )


def _org_chart_title(page_text: str) -> str:
    for line in page_text.splitlines():
        stripped = line.strip()
        if stripped:
            return stripped
    return "ORG CHART"
